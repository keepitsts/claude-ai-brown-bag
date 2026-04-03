#!/usr/bin/env python3
"""Generate the brown bag presentation — focused on the engineering process."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

# Colors
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2D, 0x3E, 0x6B)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MED = RGBColor(0x47, 0x55, 0x69)
SUBTLE = RGBColor(0x94, 0xA3, 0xB8)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=18, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def multi_txt(slide, lines, l, t, w, h, size=16, color=DARK, spacing=Pt(6),
              font="Calibri", bold=False):
    """Multiple paragraphs in one text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.font.bold = bold
        p.space_after = spacing
    return tb


def arrow_chain(slide, items, top, left_start, item_w, gap, box_color, text_color,
                arrow_color=SUBTLE, size=13):
    """Draw a chain of boxes with arrows between them."""
    for i, item in enumerate(items):
        left = left_start + (item_w + gap) * i
        rect(slide, left, top, item_w, Inches(0.45), box_color)
        txt(slide, item, left, top, item_w, Inches(0.45),
            size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)
        if i < len(items) - 1:
            arrow_left = left + item_w + Inches(0.02)
            txt(slide, "\u2192", arrow_left, top - Inches(0.02), Inches(gap - 0.04),
                Inches(0.45), size=18, color=arrow_color, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Building the System\nThat Builds the System",
    Inches(1.5), Inches(1.5), Inches(10), Inches(2),
    size=48, color=WHITE, bold=True)

rect(s, Inches(1.5), Inches(3.7), Inches(2.5), Inches(0.04), ACCENT)

txt(s, "A Software Engineer's Guide to Claude Code",
    Inches(1.5), Inches(4.1), Inches(10), Inches(0.6),
    size=22, color=SUBTLE)

txt(s, "Rudy Zauel  |  Simple Technology Solutions  |  April 3, 2026",
    Inches(1.5), Inches(5.0), Inches(8), Inches(0.5),
    size=16, color=SUBTLE)


# ============================================================
# SLIDE 2: The Question
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "The Experiment", Inches(1), Inches(0.6), Inches(11), Inches(0.7),
    size=36, color=NAVY, bold=True)

txt(s, (
    "I had a spec for a 7-feature application with 6 API integrations.\n"
    "Can Claude Code just... build it while I walk my dog?"
), Inches(1.2), Inches(1.8), Inches(10.5), Inches(1.2),
    size=22, color=DARK_BLUE)

# The naive approach
rect(s, Inches(1), Inches(3.5), Inches(5.3), Inches(3.2), LIGHT)
txt(s, "What I Wanted", Inches(1.3), Inches(3.6), Inches(4.5), Inches(0.4),
    size=18, color=NAVY, bold=True)
multi_txt(s, [
    '"Here\'s my spec. Build it."',
    "Go walk Milo (28lb Cavapoo).",
    "Come back to a working app.",
], Inches(1.5), Inches(4.2), Inches(4.5), Inches(2.0),
    size=16, color=MED)

# What actually happens
rect(s, Inches(7), Inches(3.5), Inches(5.3), Inches(3.2), NAVY)
txt(s, "What Actually Happens", Inches(7.3), Inches(3.6), Inches(4.5), Inches(0.4),
    size=18, color=RED, bold=True)
multi_txt(s, [
    "Context window fills up by feature 4.",
    "It forgets patterns from feature 1.",
    "Mistakes compound without verification.",
    "Quality degrades well before the limit.",
    "You come back to a mess.",
], Inches(7.5), Inches(4.2), Inches(4.5), Inches(2.5),
    size=16, color=SUBTLE)


# ============================================================
# SLIDE 3: The Insight
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "The insight:", Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
    size=24, color=SUBTLE)

txt(s, "I needed to build the system\nthat would build the system.",
    Inches(1), Inches(2.3), Inches(11.3), Inches(1.5),
    size=40, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

rect(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.04), ACCENT)

txt(s, (
    "Not one giant prompt. A pipeline of markdown artifacts — "
    "each one the input to the next — that refines an idea "
    "into a deployable application."
), Inches(1), Inches(4.7), Inches(11.3), Inches(1.0),
    size=20, color=SUBTLE)


# ============================================================
# SLIDE 4: The Refinement Chain
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "The Refinement Chain", Inches(1), Inches(0.6), Inches(11), Inches(0.7),
    size=36, color=NAVY, bold=True)

txt(s, "Each markdown file was the input that produced the next one",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4),
    size=18, color=MED)

chain = [
    ("Company\nBriefing", "Who is STS?\nWhat do they care about?", "research/", ACCENT),
    ("Demo Spec", "7 features, 6 APIs,\ntarget persona, scope", "research/", ACCENT),
    ("Tech Stack\nDecision Record", "Every choice with\nrationale, not just 'use React'", "research/", ACCENT),
    ("Build Process\nResearch", "What agentic harness?\nHeadless pipeline recommended", "research/", ORANGE),
    ("CLAUDE.md +\n11 Build Prompts", "Conventions + per-feature\nspec files", "build/prompts/", GREEN),
    ("Working\nApplication", "4,455 lines, 58 files,\n7 features, deployed", "trailhead/", GREEN),
]

for i, (title, desc, loc, color) in enumerate(chain):
    left = Inches(0.5) + Inches(2.1) * i
    top = Inches(2.2)

    # Card
    rect(s, left, top, Inches(1.9), Inches(2.8), LIGHT)
    rect(s, left, top, Inches(1.9), Inches(0.06), color)

    # Number
    txt(s, str(i + 1), left + Inches(0.15), top + Inches(0.2),
        Inches(0.4), Inches(0.4), size=20, color=color, bold=True)

    # Title
    txt(s, title, left + Inches(0.15), top + Inches(0.6),
        Inches(1.6), Inches(0.7), size=14, color=NAVY, bold=True)

    # Description
    txt(s, desc, left + Inches(0.15), top + Inches(1.5),
        Inches(1.6), Inches(0.8), size=11, color=MED)

    # Location
    txt(s, loc, left + Inches(0.15), top + Inches(2.4),
        Inches(1.6), Inches(0.3), size=10, color=SUBTLE, font="Consolas")

    # Arrow between cards
    if i < len(chain) - 1:
        arrow_l = left + Inches(1.95)
        txt(s, "\u2192", arrow_l, top + Inches(1.0), Inches(0.2), Inches(0.5),
            size=22, color=SUBTLE, bold=True, align=PP_ALIGN.CENTER)

# Bottom insight
rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), NAVY)
txt(s, "Every artifact was a markdown file.  The entire process is in the repo.",
    Inches(1), Inches(5.7), Inches(11.3), Inches(0.5),
    size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "prompts/001 through prompts/017  |  research/*.md  |  build/prompts/00-10",
    Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
    size=14, color=SUBTLE, align=PP_ALIGN.CENTER, font="Consolas")


# ============================================================
# SLIDE 5: Step 1 — Research First
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Step 1:  Research First, Code Later", Inches(1), Inches(0.6),
    Inches(11), Inches(0.7), size=36, color=NAVY, bold=True)

txt(s, "I used claude.ai (with web search) for research, Claude Code for implementation",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4), size=18, color=MED)

# Three research artifacts
artifacts = [
    ("Company Briefing",
     "simple-technology-solutions-a-pre-presentation-briefing...md",
     [
         "Who is STS? $60M DOI contract, 3,000 officers, 855+ locations",
         "Their AI products: C-CAT, LEI, A\u00b3",
         "Culture: mission-first, 'Keep IT Simple'",
         "Key people and what they care about",
     ]),
    ("Demo Spec",
     "sts-demo-spec.md",
     [
         "7 features defined in detail (F1\u2013F7)",
         "6 free APIs identified with auth requirements",
         "Target persona: DOI law enforcement officer",
         "What the demo IS and IS NOT",
         "Success criteria for the presentation",
     ]),
    ("Tech Stack Decision Record",
     "trailhead-tech-stack.md",
     [
         "Every choice documented with rationale",
         "Not 'use React' but 'why Next.js App Router'",
         "State management: Zustand + TanStack Query and why",
         "What was deliberately left out (database, auth, tests)",
         "Full dependency list and project structure",
     ]),
]

for i, (title, filename, items) in enumerate(artifacts):
    left = Inches(0.7) + Inches(4.1) * i
    top = Inches(2.2)

    rect(s, left, top, Inches(3.8), Inches(4.5), LIGHT)
    rect(s, left, top, Inches(3.8), Inches(0.06), ACCENT)

    txt(s, title, left + Inches(0.3), top + Inches(0.2),
        Inches(3.2), Inches(0.4), size=20, color=NAVY, bold=True)

    txt(s, filename, left + Inches(0.3), top + Inches(0.65),
        Inches(3.2), Inches(0.3), size=10, color=SUBTLE, font="Consolas")

    for j, item in enumerate(items):
        txt(s, f"\u2022  {item}", left + Inches(0.3), top + Inches(1.1) + Inches(0.42) * j,
            Inches(3.2), Inches(0.4), size=13, color=MED)


# ============================================================
# SLIDE 6: Step 2 — Research the Build Process
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, 'Step 2:  "I Need an Agentic Harness"', Inches(1), Inches(0.6),
    Inches(11), Inches(0.7), size=36, color=NAVY, bold=True)

txt(s, "Before writing any code, I researched how to make Claude Code build reliably",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4), size=18, color=MED)

# The research prompt excerpt
rect(s, Inches(0.8), Inches(2.0), Inches(7.2), Inches(4.8), LIGHT)
txt(s, "The Research Prompt I Wrote", Inches(1.1), Inches(2.1),
    Inches(6), Inches(0.4), size=16, color=NAVY, bold=True)

prompt_lines = [
    "What I want:  give the spec to Claude Code, walk my dog, come back to a working app.",
    "",
    "What I know won't work:  Claude Code won't build 7 features in one prompt.",
    "Context fills up. It forgets patterns. Mistakes compound.",
    "",
    "What I need:  an approach that decomposes the spec into right-sized units,",
    "sequences the build, verifies each unit, and minimizes human intervention.",
    "",
    "Evaluate:  single session, headless orchestration, Agent SDK pipeline,",
    "/batch, agent teams, custom subagents, or something I haven't thought of.",
    "",
    "Be opinionated \u2014 I want 'here's what you should do' not '5 options to consider.'",
]
multi_txt(s, prompt_lines, Inches(1.1), Inches(2.6), Inches(6.8), Inches(4.0),
    size=13, color=MED, font="Consolas", spacing=Pt(3))

# What the research recommended
rect(s, Inches(8.3), Inches(2.0), Inches(4.3), Inches(4.8), NAVY)
txt(s, "The Recommendation", Inches(8.6), Inches(2.1),
    Inches(3.7), Inches(0.4), size=16, color=ACCENT, bold=True)

rec_items = [
    ("Sequential headless pipeline", "claude -p with fresh context per feature"),
    ("TypeScript verification gates", "tsc --noEmit + build between every step"),
    ("Auto-commit on success", "Clean git history, one commit per feature"),
    ("CLAUDE.md under 80 lines", "Persistent memory across fresh contexts"),
    ("11 self-contained prompts", "Each prompt file is a complete feature spec"),
]

for i, (title, desc) in enumerate(rec_items):
    top = Inches(2.7) + Inches(0.75) * i
    txt(s, title, Inches(8.6), top, Inches(3.7), Inches(0.35),
        size=15, color=WHITE, bold=True)
    txt(s, desc, Inches(8.6), top + Inches(0.3), Inches(3.7), Inches(0.35),
        size=12, color=SUBTLE)

# Rejected options
txt(s, "Rejected: Agent Teams (experimental, unstable), /batch (can't handle shared\n"
       "dependencies), single long session (context degrades by feature 4)",
    Inches(8.6), Inches(6.0), Inches(3.7), Inches(0.7),
    size=11, color=SUBTLE)


# ============================================================
# SLIDE 7: Step 3 — CLAUDE.md
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Step 3:  CLAUDE.md \u2014 50 Lines of Persistent Memory",
    Inches(1), Inches(0.6), Inches(11), Inches(0.7),
    size=36, color=NAVY, bold=True)

txt(s, "Every fresh claude -p call loads this file. It's how conventions survive across context windows.",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4), size=18, color=MED)

# Show the actual CLAUDE.md content
rect(s, Inches(0.8), Inches(2.0), Inches(7.5), Inches(5.0), RGBColor(0x0D, 0x11, 0x17))

claude_md = [
    "# TRAILHEAD \u2014 Federal Lands Daily Briefing Tool",
    "",
    "## Tech Stack",
    "- Next.js 15 App Router, TypeScript strict",
    "- Tailwind CSS v4 (CSS-based config, NOT tailwind.config.js)",
    "- shadcn/ui, Zustand for client state, TanStack Query v5",
    "",
    "## Critical Conventions",
    "- Server Components by default. 'use client' ONLY for interactivity",
    "- Next.js 15: params are Promises \u2014 must await",
    "- React Leaflet REQUIRES dynamic import with ssr: false",
    "- Data flow: lib/api/ \u2192 app/api/ route \u2192 TanStack Query hook \u2192 component",
    "- Named exports. Conventional commits (feat:, fix:, refactor:)",
    "",
    "## Verification (run after every feature)",
    "npx tsc --noEmit && npm run build",
]
multi_txt(s, claude_md, Inches(1.1), Inches(2.2), Inches(7.0), Inches(4.5),
    size=13, color=GREEN, font="Consolas", spacing=Pt(2))

# Why it matters
rect(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(5.0), LIGHT)
txt(s, "Why This Matters", Inches(9.1), Inches(2.2), Inches(3.2), Inches(0.4),
    size=18, color=NAVY, bold=True)

insights = [
    "Each feature builds in a fresh context window \u2014 no memory of previous features",
    "CLAUDE.md is the only thing that carries forward",
    "It tells Claude the patterns, not the history",
    "Under 80 lines \u2014 every line competes for the model's attention",
    "Claude rediscovers the codebase from the filesystem, guided by these conventions",
    "The data flow pattern (lib/api \u2192 route \u2192 hook \u2192 component) kept all 11 features consistent",
]
for i, insight in enumerate(insights):
    top = Inches(2.9) + Inches(0.6) * i
    rect(s, Inches(9.1), top, Inches(0.06), Inches(0.4), ACCENT)
    txt(s, insight, Inches(9.4), top, Inches(3.0), Inches(0.5),
        size=12, color=MED)


# ============================================================
# SLIDE 8: Step 4 — The Build Pipeline
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Step 4:  The Build Pipeline", Inches(1), Inches(0.6),
    Inches(11), Inches(0.7), size=36, color=NAVY, bold=True)

txt(s, "A 30-line bash script that runs 11 headless Claude Code calls",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4), size=18, color=MED)

# build.sh pseudocode
rect(s, Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.8), RGBColor(0x0D, 0x11, 0x17))

build_code = [
    "#!/bin/bash",
    "",
    "build_step() {",
    "  local step_name=\"$1\"",
    "  local prompt=\"$2\"",
    "",
    "  # Fresh context window for each feature",
    "  claude -p \"$prompt\" \\",
    '    --allowedTools "Read,Write,Edit,Bash,Glob,Grep"',
    "",
    "  # Verification gate",
    "  npx tsc --noEmit",
    "  pnpm build",
    "",
    "  # Only commit on success",
    '  git add -A && git commit -m "feat: $step_name"',
    "}",
    "",
    "# Run all 11 features sequentially",
    'build_step "scaffold"    "$(cat build/prompts/00-scaffold.md)"',
    'build_step "shared-infra" "$(cat build/prompts/01-shared-infra.md)"',
    'build_step "base-layout" "$(cat build/prompts/02-base-layout.md)"',
    "# ... through 10-integration.md",
]
multi_txt(s, build_code, Inches(1.0), Inches(2.2), Inches(6.0), Inches(4.5),
    size=12, color=GREEN, font="Consolas", spacing=Pt(1))

# The 11 steps with results
rect(s, Inches(7.8), Inches(2.0), Inches(4.8), Inches(4.8), LIGHT)
txt(s, "11 Steps, 0 Failures", Inches(8.1), Inches(2.1),
    Inches(4.2), Inches(0.4), size=18, color=NAVY, bold=True)

steps = [
    ("00", "Scaffold", "Next.js 15 + Tailwind v4 + shadcn/ui"),
    ("01", "Shared Infra", "Types, API clients, stores, locations"),
    ("02", "Base Layout", "Sidebar, location selector, dashboard shell"),
    ("03", "Weather", "Open-Meteo + 15min auto-refresh"),
    ("04", "Seismic", "USGS earthquakes + 5min refresh"),
    ("05", "Fires", "NASA FIRMS + CSV parsing"),
    ("06", "Crime", "FBI Crime Data + rate aggregation"),
    ("07", "Astronomy", "Sunrise/sunset/moon phase"),
    ("08", "Map", "React Leaflet + fire/earthquake markers"),
    ("09", "Briefing", "Claude API streaming + 5-section format"),
    ("10", "Integration", "Watch list, responsive, error boundaries"),
]

for i, (num, name, detail) in enumerate(steps):
    top = Inches(2.65) + Inches(0.37) * i
    txt(s, num, Inches(8.1), top, Inches(0.4), Inches(0.3),
        size=11, color=ACCENT, bold=True, font="Consolas")
    txt(s, name, Inches(8.5), top, Inches(1.5), Inches(0.3),
        size=12, color=NAVY, bold=True)
    txt(s, detail, Inches(10.0), top, Inches(2.4), Inches(0.3),
        size=11, color=MED)


# ============================================================
# SLIDE 9: What a Build Prompt Looks Like
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "What a Build Prompt Looks Like", Inches(1), Inches(0.6),
    Inches(11), Inches(0.7), size=36, color=NAVY, bold=True)

txt(s, "Each of the 11 prompts is a self-contained feature spec (~200-400 words)",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4), size=18, color=MED)

# Example prompt structure
rect(s, Inches(0.8), Inches(2.0), Inches(7.5), Inches(5.0), RGBColor(0x0D, 0x11, 0x17))

prompt_example = [
    "# Feature: USGS Earthquake Monitoring",
    "",
    "## What to build",
    "- API route at app/api/seismic/route.ts",
    "- Earthquake list component with magnitude-coded badges",
    "- Map overlay with circle markers sized by magnitude",
    "- TanStack Query hook in hooks/useSeismic.ts",
    "",
    "## API Details",
    "- Endpoint: https://earthquake.usgs.gov/fdsnws/event/1/query",
    "- Parameters: format=geojson, latitude, longitude, maxradiuskm",
    "- No authentication required",
    "",
    "## Patterns to follow     \u2190 THIS IS THE KEY SECTION",
    "- See lib/api/weather.ts for API client wrapper pattern",
    "- See app/api/weather/route.ts for route handler pattern",
    "- See hooks/useWeather.ts for TanStack Query hook pattern",
    "",
    "## Scope Control",
    "- Implement ONLY what is specified above",
    "- Do not add features beyond these requirements",
]
multi_txt(s, prompt_example, Inches(1.0), Inches(2.2), Inches(7.0), Inches(4.6),
    size=13, color=GREEN, font="Consolas", spacing=Pt(1))

# Annotations
rect(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(5.0), LIGHT)
txt(s, "Key Elements", Inches(9.1), Inches(2.2), Inches(3.2), Inches(0.4),
    size=18, color=NAVY, bold=True)

elements = [
    ('"Patterns to follow"', "Points Claude to existing code.\nIt reads those files first,\nthen follows the same structure."),
    ('"API Details"', "Exact endpoints, parameters,\nauth requirements. No guessing."),
    ('"What to build"', "Concrete deliverables.\nFiles that should exist when done."),
    ('"Scope Control"', "Prevents Claude from adding\nauth, dark mode, i18n, etc."),
]

for i, (title, desc) in enumerate(elements):
    top = Inches(2.9) + Inches(1.0) * i
    txt(s, title, Inches(9.1), top, Inches(3.2), Inches(0.35),
        size=14, color=ACCENT, bold=True)
    txt(s, desc, Inches(9.1), top + Inches(0.35), Inches(3.2), Inches(0.6),
        size=12, color=MED)


# ============================================================
# SLIDE 10: The Debug Loop
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Step 5:  The Debug Loop", Inches(1), Inches(0.6),
    Inches(11), Inches(0.7), size=36, color=NAVY, bold=True)

txt(s, "The pipeline built the app. Then reality hit \u2014 3 out of 6 APIs had issues.",
    Inches(1), Inches(1.3), Inches(11), Inches(0.4), size=18, color=MED)

# Debug method
rect(s, Inches(0.8), Inches(2.1), Inches(12.0), Inches(1.0), LIGHT)
arrow_chain(s,
    ["Open app in browser", "See error", "Save console log\nto debug-logs/", "Tell Claude Code\n\"read this log\"", "Claude diagnoses\nand fixes"],
    Inches(2.25), Inches(1.0), Inches(2.1), Inches(0.3), NAVY, WHITE, SUBTLE, 11)

# The four bugs
bugs = [
    ("NASA FIRMS", "400 Error",
     "URL used lat,lng,radius format \u2014 FIRMS expects bounding box west,south,east,north",
     "One prompt, fixed"),
    ("FBI Crime Data", "503 + 404",
     "Env var mismatch (FBI_API_KEY vs FBI_CRIME_API_KEY) AND wrong endpoint format entirely.\n"
     "Had to rewrite to /summarized/state/{state}/{offense}?from=MM-YYYY&to=MM-YYYY",
     "Two prompts, complete rewrite"),
    ("Weather", "Wrong units",
     "Open-Meteo defaults to Celsius/km/h. Needed temperature_unit=fahrenheit, wind_speed_unit=mph",
     "One line fix"),
    ("AI Briefing", "Empty response",
     "Anthropic API key had zero credits. Not a code bug at all.",
     "No code fix \u2014 bought credits"),
]

for i, (api, error, cause, resolution) in enumerate(bugs):
    top = Inches(3.5) + Inches(0.95) * i
    rect(s, Inches(0.8), top, Inches(0.06), Inches(0.75), RED)
    txt(s, f"{api}  \u2014  {error}", Inches(1.1), top,
        Inches(3), Inches(0.35), size=15, color=NAVY, bold=True)
    txt(s, cause, Inches(4.3), top, Inches(5.5), Inches(0.7),
        size=12, color=MED)
    txt(s, resolution, Inches(10.2), top, Inches(2.5), Inches(0.7),
        size=12, color=GREEN, bold=True)


# ============================================================
# SLIDE 11: The Result
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "The Result", Inches(1), Inches(0.5), Inches(11), Inches(0.6),
    size=36, color=WHITE, bold=True)

# Stats
stats = [
    ("4,455", "Lines of TypeScript"),
    ("58", "Source Files"),
    ("20", "Git Commits"),
    ("11", "Build Steps"),
    ("0", "Build Failures"),
]
for i, (num, label) in enumerate(stats):
    left = Inches(0.6) + Inches(2.55) * i
    rect(s, left, Inches(1.4), Inches(2.3), Inches(1.5), DARK_BLUE)
    nc = GREEN if num == "0" else ACCENT
    txt(s, num, left, Inches(1.5), Inches(2.3), Inches(0.8),
        size=44, color=nc, bold=True, align=PP_ALIGN.CENTER)
    txt(s, label, left, Inches(2.25), Inches(2.3), Inches(0.4),
        size=13, color=SUBTLE, align=PP_ALIGN.CENTER)

# Live URL
rect(s, Inches(1), Inches(3.4), Inches(11.3), Inches(0.6), ACCENT)
txt(s, "trailhead-pearl.vercel.app", Inches(1), Inches(3.4),
    Inches(11.3), Inches(0.6), size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# What it does
txt(s, "TRAILHEAD \u2014 Federal Lands Daily Briefing Tool",
    Inches(1), Inches(4.4), Inches(11.3), Inches(0.5),
    size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

features = [
    "Real-time weather, seismic, wildfire, crime, and astronomy data for 25 federal land sites",
    "AI-generated operational briefings via Claude API with streaming",
    "Interactive Leaflet map with fire and earthquake overlays",
    "Briefing history, multi-location watchlist, and shareable export",
    "Deployed to Vercel with all API keys configured",
]
for i, f in enumerate(features):
    txt(s, f"\u2022  {f}", Inches(2), Inches(5.1) + Inches(0.38) * i,
        Inches(9.3), Inches(0.35), size=15, color=SUBTLE)


# ============================================================
# SLIDE 12: Live Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Live Demo", Inches(1), Inches(2.5), Inches(11.3), Inches(1),
    size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, "trailhead-pearl.vercel.app", Inches(1), Inches(3.7),
    Inches(11.3), Inches(0.6), size=24, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13: Lessons Learned
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "What I Learned", Inches(1), Inches(0.6), Inches(11), Inches(0.7),
    size=36, color=NAVY, bold=True)

lessons = [
    ("Don't prompt. Engineer.",
     "The prompt is the last step. Before it, you need research, a spec, a tech stack "
     "decision record, a build pipeline, and a conventions file. The system that builds "
     "the system is the real work."),
    ("Fresh context beats long context.",
     "11 separate claude -p calls with 50 lines of CLAUDE.md outperformed one long session "
     "that degraded by feature 4. Claude rediscovers the codebase from the filesystem \u2014 "
     "you just need to tell it the patterns."),
    ("Verification gates are non-negotiable.",
     "TypeScript compilation + build between every step caught errors before they compounded. "
     "Without this, feature 7 would inherit broken assumptions from feature 3."),
    ("Markdown in, software out.",
     "The entire process was: write a markdown file, feed it to Claude, get an artifact, "
     "use that artifact to write the next markdown file. The repo is the complete record."),
    ("AI gets you 0\u201380%. You're still the engineer.",
     "API format mismatches, env var typos, wrong units \u2014 real-world integration still "
     "needs a human who reads the error and understands the domain."),
]

for i, (title, desc) in enumerate(lessons):
    top = Inches(1.6) + Inches(1.1) * i
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), top + Inches(0.05),
                                 Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    txt(s, str(i + 1), Inches(1), top + Inches(0.05), Inches(0.45), Inches(0.45),
        size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.7), top, Inches(10), Inches(0.4),
        size=18, color=NAVY, bold=True)
    txt(s, desc, Inches(1.7), top + Inches(0.4), Inches(10), Inches(0.55),
        size=14, color=MED)


# ============================================================
# SLIDE 14: Q&A
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
rect(s, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

txt(s, "Questions?", Inches(1), Inches(2.0), Inches(11.3), Inches(1),
    size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), ACCENT)

txt(s, "App:   trailhead-pearl.vercel.app",
    Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
    size=18, color=SUBTLE, align=PP_ALIGN.CENTER)
txt(s, "Repo:  github.com/keepitsts/claude-ai-brown-bag",
    Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
    size=18, color=SUBTLE, align=PP_ALIGN.CENTER)

rect(s, Inches(2), Inches(5.3), Inches(9.3), Inches(1.2), DARK_BLUE)
txt(s, "Every prompt, every research artifact, every build script\nis in the repo. "
       "The entire process is transparent.",
    Inches(2.5), Inches(5.5), Inches(8.3), Inches(0.8),
    size=16, color=SUBTLE, align=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output = "/home/rczauel/claude-ai-brown-bag/presentation/claude-code-brown-bag.pptx"
prs.save(output)
print(f"Saved: {output}")
